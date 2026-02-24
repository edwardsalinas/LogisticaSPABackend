from pptx import Presentation
import os

def extract_text_from_ppt(filename):
    print(f"Extracting text from {filename}...")
    try:
        prs = Presentation(filename)
        for i, slide in enumerate(prs.slides):
            print(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    print(shape.text)
    except Exception as e:
        print(f"Error extracting text: {e}")

if __name__ == "__main__":
    # Note: python-pptx works for .pptx, might fail for old .ppt
    extract_text_from_ppt("Copia de SMART Goals.ppt")
